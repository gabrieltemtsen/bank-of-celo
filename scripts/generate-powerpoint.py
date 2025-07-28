#!/usr/bin/env python3
"""
Generate PowerPoint presentation for Bank of Celo Proof of Ship Season 6 submission
"""

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    import os
except ImportError:
    print("Installing required packages...")
    import subprocess
    import sys
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE

def create_bank_of_celo_presentation():
    # Create presentation
    prs = Presentation()
    
    # Define colors
    celo_green = RGBColor(52, 208, 127)  # #34d07f
    base_purple = RGBColor(139, 92, 246)  # #8b5cf6
    dark_blue = RGBColor(51, 65, 85)      # #334155
    light_gray = RGBColor(248, 250, 252)  # #f8fafc
    
    # Slide 1: Title Slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add bank emoji as logo
    logo = slide1.shapes.add_textbox(Inches(4.5), Inches(0.5), Inches(1), Inches(1))
    logo_frame = logo.text_frame
    logo_p = logo_frame.paragraphs[0]
    logo_p.text = "🏦"
    logo_p.font.size = Pt(72)
    logo_p.alignment = PP_ALIGN.CENTER
    
    # Title
    title = slide1.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
    title_frame = title.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = "Bank of Celo"
    title_p.font.size = Pt(48)
    title_p.font.bold = True
    title_p.font.color.rgb = celo_green
    title_p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle = slide1.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(1))
    subtitle_frame = subtitle.text_frame
    subtitle_p = subtitle_frame.paragraphs[0]
    subtitle_p.text = "Social DeFi Banking Platform"
    subtitle_p.font.size = Pt(28)
    subtitle_p.alignment = PP_ALIGN.CENTER
    
    # Features
    features = slide1.shapes.add_textbox(Inches(2), Inches(4.5), Inches(6), Inches(1))
    features_frame = features.text_frame
    features_p = features_frame.paragraphs[0]
    features_p.text = "Multi-Chain • Farcaster Native • Community Driven"
    features_p.font.size = Pt(20)
    features_p.alignment = PP_ALIGN.CENTER
    
    # Proof of Ship badge
    pos_badge = slide1.shapes.add_textbox(Inches(2.5), Inches(6), Inches(5), Inches(1))
    pos_frame = pos_badge.text_frame
    pos_p = pos_frame.paragraphs[0]
    pos_p.text = "Proof of Ship Season 6 Submission"
    pos_p.font.size = Pt(24)
    pos_p.font.bold = True
    pos_p.font.color.rgb = RGBColor(251, 191, 36)  # Golden
    pos_p.alignment = PP_ALIGN.CENTER
    
    # Slide 2: Project Overview
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    slide2.shapes.title.text = "🌟 Project Overview"
    
    content = slide2.shapes.placeholders[1].text_frame
    content.text = """🏦 Social Banking
Community-driven vault system with transparent donations and automated rewards

🔗 Multi-Chain
Seamless operation across Celo and Base networks with unified UX

🎯 Farcaster Native
Built specifically for the Farcaster ecosystem with Frame SDK integration

💰 Yield Generation
Multiple earning opportunities: savings, rewards, jackpots, and check-ins"""
    
    # Slide 3: Core Features
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "🚀 Core Features"
    
    content3 = slide3.shapes.placeholders[1].text_frame
    content3.text = """Banking System:
• Donation Vaults with tiered rewards
• Claim System with Farcaster verification
• Gasless Transactions with EIP-712
• Quality Score integration (>0.39)

Gamification:
• Jackpot/Lottery system
• Daily Check-ins for rewards
• Leaderboards (Top 100 donors)
• Tiered System (3 donor levels)

Fx Savings:
• USDC Vault on Base → Earn DEGEN
• cEUR Vault on Celo → Earn CELO
• Real-time APY calculation
• Token Swapping via Farcaster SDK"""
    
    # Slide 4: Multi-Chain Architecture
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "🌐 Multi-Chain Architecture"
    
    content4 = slide4.shapes.placeholders[1].text_frame
    content4.text = """Celo Network - Native CELO Operations:
• Bank Vault: Native CELO deposits
• Claim Amount: 0.5 CELO
• Jackpot Tickets: 1 CELO each
• Fx Savings: cEUR → CELO rewards
• 8 Deployed Contracts

Base Network - DEGEN Token Operations:
• DEGEN Bank: ERC-20 token support
• Claim Amount: 100 DEGEN
• Jackpot Tickets: 250 DEGEN each
• Fx Savings: USDC → DEGEN rewards
• 4 Deployed Contracts

✨ Unified User Experience Across Both Networks"""
    
    # Slide 5: Farcaster Integration
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "🎭 Deep Farcaster Integration"
    
    content5 = slide5.shapes.placeholders[1].text_frame
    content5.text = """Identity Verification:
• Farcaster ID Required for all claims
• Quality Score ≥0.39 minimum
• Username Resolution via Neynar API
• FID-based Blacklisting system

Frame SDK Integration:
• Native Token Swapping in frames
• CAIP-19 Standard compliance
• ETH→USDC on Base
• CELO→cEUR on Celo

Social Features:
• Leaderboard Display with usernames
• Community Recognition system
• Social-first DeFi experience
• Engagement Tracking"""
    
    # Slide 6: Technical Excellence
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    slide6.shapes.title.text = "⚡ Technical Highlights"
    
    content6 = slide6.shapes.placeholders[1].text_frame
    content6.text = """Architecture:
• Next.js 15 with TypeScript
• Wagmi v2 for Web3 integration
• Framer Motion animations
• TailwindCSS with dark mode
• Convex real-time data

Security:
• EIP-712 typed signatures
• ReentrancyGuard protection
• Input Validation & sanitization
• Rate Limiting mechanisms
• Multi-layer access controls

Performance:
• Real-time Sync every 3 seconds
• Optimistic UI updates
• Contract Read caching
• Mobile-first responsive design"""
    
    # Slide 7: Smart Contract Innovation
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    slide7.shapes.title.text = "📜 Smart Contract Innovation"
    
    content7 = slide7.shapes.placeholders[1].text_frame
    content7.text = """Key Metrics:
• 14 Smart Contracts deployed
• 2 Networks supported
• 100 Leaderboard slots
• 3% Jackpot win rate

Banking Contracts:
• Donation Vaults: Secure CELO/DEGEN storage
• Claim System: Farcaster ID verification
• Tier Management: Automated donor classification
• Fee Distribution: 10% dev fee structure

Fx Vaults:
• USDC Vault: Base network USDC deposits
• cEUR Vault: Celo network cEUR deposits
• Reward Distribution: Automated yield payments
• Real-time APY: Dynamic rate calculation"""
    
    # Slide 8: Economic Model
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    slide8.shapes.title.text = "💎 Economic Model & Rewards"
    
    content8 = slide8.shapes.placeholders[1].text_frame
    content8.text = """Vault Claims:
• Celo: 0.5 CELO per claim
• Base: 100 DEGEN per claim
• Cooldown: Quality score based
• Verification: Farcaster ID required

Fx Savings:
• Dynamic APY: Real-time rate calculation
• Dual Networks: USDC & cEUR support
• Flexible Deposits: No lock-up periods
• Compound Rewards: Automatic reinvestment

Jackpot Lottery:
• Ticket Price: 1 CELO / 250 DEGEN
• Win Probability: 3% base rate
• Dev Fee: 5% on winnings
• Carryover: Unclaimed prizes accumulate

🔥 10x Multiplier for Identity-Verified Users"""
    
    # Slide 9: User Journey
    slide9 = prs.slides.add_slide(prs.slide_layouts[1])
    slide9.shapes.title.text = "🎯 User Experience Journey"
    
    content9 = slide9.shapes.placeholders[1].text_frame
    content9.text = """1️⃣ Onboarding:
• Connect Farcaster account
• Verify quality score
• Choose network (Celo/Base)
• Complete welcome tutorial

2️⃣ Engagement:
• Daily check-ins for rewards
• Donate to vault system
• Participate in jackpots
• Climb leaderboards

3️⃣ Earning:
• Claim vault rewards
• Earn from Fx savings
• Win jackpot prizes
• Receive tier bonuses

Key Differentiators: Farcaster Native • Gasless Claims • Multi-Chain • Social DeFi"""
    
    # Slide 10: Innovation Impact
    slide10 = prs.slides.add_slide(prs.slide_layouts[1])
    slide10.shapes.title.text = "🚀 Innovation & Impact"
    
    content10 = slide10.shapes.placeholders[1].text_frame
    content10.text = """Farcaster Ecosystem Impact:
• 50M+ Funding ecosystem opportunity
• Social-first DeFi user experience
• Frame SDK native integration
• Community Building through finance

Technical Innovation:
• Cross-chain UX unified experience
• EIP-712 Gasless transactions
• Quality Score access control
• Real-time Sync performance

Market Opportunity:
• Web2 → Web3 bridge for social users
• Emerging Markets financial inclusion
• Creator Economy monetization
• Community Finance new paradigm

🏆 First truly social-native DeFi platform combining Farcaster identity, 
multi-chain infrastructure, and gamified banking"""
    
    # Slide 11: Key Metrics
    slide11 = prs.slides.add_slide(prs.slide_layouts[1])
    slide11.shapes.title.text = "📊 Key Metrics & Achievements"
    
    content11 = slide11.shapes.placeholders[1].text_frame
    content11.text = """Core Metrics:
• 14 Smart Contracts
• 2 Networks
• 100 Leaderboard Spots
• 5 Core Features
• 3% Jackpot Win Rate
• 10x Verified User Bonus
• 0.39 Min Quality Score

Celo Network Contracts:
• Bank Vault: 0x18Ea8d1D41A3307D159D2d3C1fCfBCF139354A8F
• Jackpot: 0x9602d02Bd17d9f1c1EB09028fCea26dD29383611
• cEUR Vault: 0x6C617A05b9D183D2BD1A3350F4782Fc125460634

Base Network Contracts:
• USDC Vault: 0x8Ca054b89C9F04C5546f37B633690fb940Cc4130
• DEGEN Bank: Multi-contract system"""
    
    # Slide 12: Conclusion
    slide12 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Main title
    conclusion_title = slide12.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1.5))
    conclusion_title_frame = conclusion_title.text_frame
    conclusion_title_p = conclusion_title_frame.paragraphs[0]
    conclusion_title_p.text = "🌟 The Future of Social DeFi"
    conclusion_title_p.font.size = Pt(36)
    conclusion_title_p.font.bold = True
    conclusion_title_p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    conclusion_subtitle = slide12.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    conclusion_subtitle_frame = conclusion_subtitle.text_frame
    conclusion_subtitle_p = conclusion_subtitle_frame.paragraphs[0]
    conclusion_subtitle_p.text = "Bank of Celo: Where Social Meets Finance"
    conclusion_subtitle_p.font.size = Pt(24)
    conclusion_subtitle_p.font.color.rgb = celo_green
    conclusion_subtitle_p.alignment = PP_ALIGN.CENTER
    
    # Content
    conclusion_content = slide12.shapes.add_textbox(Inches(1.5), Inches(3.5), Inches(7), Inches(2.5))
    conclusion_content_frame = conclusion_content.text_frame
    conclusion_content_p = conclusion_content_frame.paragraphs[0]
    conclusion_content_p.text = """Mission: Bridge Web2 social users into Web3 finance through familiar, gamified experiences

Vision: Become the standard for social DeFi, empowering communities to build financial systems

Innovation: Pioneering the convergence of social graphs and financial infrastructure

Impact: Creating financial inclusion through social participation"""
    conclusion_content_p.font.size = Pt(16)
    
    # Final call to action
    cta = slide12.shapes.add_textbox(Inches(2), Inches(6), Inches(6), Inches(1.5))
    cta_frame = cta.text_frame
    cta_p = cta_frame.paragraphs[0]
    cta_p.text = "🏆 Ready for Proof of Ship Season 6\n\nA complete, innovative DeFi platform that advances the entire Farcaster ecosystem"
    cta_p.font.size = Pt(18)
    cta_p.font.bold = True
    cta_p.font.color.rgb = RGBColor(251, 191, 36)  # Golden
    cta_p.alignment = PP_ALIGN.CENTER
    
    # Save the presentation
    output_path = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), 
                              'public', 'Bank-of-Celo-Proof-of-Ship-S6.pptx')
    prs.save(output_path)
    print(f"PowerPoint presentation saved to: {output_path}")
    return output_path

if __name__ == "__main__":
    create_bank_of_celo_presentation()